"""Дизайн-токены и хелперы PPTX-альбома — стиль «Минимал · Спецификация».

Цвета — RGB-кортежи. Хелперы рисуют на slide элементы единого стиля:
шапку с амбер-чертой, KPI-карточки, таблицы со статусом, плашки, подвал.
"""
from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

# ─── Палитра «Спецификация» ───
INK = (0x1A, 0x1A, 0x1A)
MUTED = (0x8A, 0x8A, 0x8A)
SOFT = (0xBB, 0xBB, 0xBB)
HAIR = (0xED, 0xED, 0xED)
AMBER = (0xF5, 0xA6, 0x23)
OK = (0x15, 0x80, 0x3D)
WARN = (0xB8, 0x76, 0x00)
BAD = (0xC0, 0x39, 0x2C)
WHITE = (0xFF, 0xFF, 0xFF)
ZEBRA = (0xFA, 0xFA, 0xFA)

EMU = 914400
SW = int(13.333 * EMU)
SH = int(7.5 * EMU)

# Шрифт — как в модели/на сайте. Заголовки: тонкий Light + жирный Black
# (как «Модель застройки **территории**»); таблицы/текст — обычный Segoe UI.
FONT = "Segoe UI"
FONT_LIGHT = "Segoe UI Light"
FONT_BLACK = "Segoe UI Black"

# Отступы: слева больше (сторона брошюровки), справа — место под вертикальный
# указатель варианта.
LEFT = 1.0
CONTENT_W = 11.55
TAB_X = 12.68          # x вертикального указателя справа

CHART_COLORS = {
    "housing": "#4A4A4A", "kindergarten": "#F5A623", "school": "#C0392C",
    "parking": "#8A8A8A", "driveways": "#B8B8B8", "znop": "#2E7D32",
    "sport": "#7CB342", "engineering": "#BDBDBD", "reserve": "#E0E0E0",
    "ink": "#1A1A1A", "amber": "#F5A623", "ok": "#15803D", "bad": "#C0392C",
    "hair": "#EDEDED", "muted": "#8A8A8A",
}


def _c(rgb):
    return RGBColor(*rgb)


def _font(f):
    """Проставить шрифт модели (Segoe UI) на font-объект. Возвращает его же."""
    f.name = FONT
    return f


class Deck:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = SW
        self.prs.slide_height = SH
        self._blank = self.prs.slide_layouts[6]

    def slide(self):
        return self.prs.slides.add_slide(self._blank)

    def save(self, path: str) -> str:
        self.prs.save(path)
        return path


def _no_shadow(shape) -> None:
    try:
        shape.shadow.inherit = False
    except Exception:  # noqa: BLE001
        pass


def rect(slide, left, top, w, h, fill, line=None, line_w=0.75):
    sp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        int(left * EMU), int(top * EMU), int(w * EMU), int(h * EMU),
    )
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = _c(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = _c(line)
        sp.line.width = Pt(line_w)
    _no_shadow(sp)
    return sp


def text(slide, left, top, w, h, lines, size=13, bold=False, color=INK,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0,
         font_name=None):
    tb = slide.shapes.add_textbox(int(left * EMU), int(top * EMU),
                                  int(w * EMU), int(h * EMU))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        if isinstance(ln, tuple):
            txt, col, bd = (ln + (color, bold))[:3]
        else:
            txt, col, bd = ln, color, bold
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.alignment = align
        p.line_spacing = spacing
        p.font.size = Pt(size)
        p.font.bold = bd
        p.font.color.rgb = _c(col)
        p.font.name = font_name or FONT
    return tb


def title_band(slide, title_plain, title_bold="", idx=None):
    """Шапка слайда: тонкий (Segoe UI Light) заголовок + жирное (Black)
    ключевое слово, по ПРАВОМУ краю; амбер-черта справа под ним."""
    right = LEFT + CONTENT_W
    tb = slide.shapes.add_textbox(int(LEFT * EMU), int(0.34 * EMU),
                                  int(CONTENT_W * EMU), int(0.75 * EMU))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r1 = p.add_run()
    r1.text = title_plain
    r1.font.size = Pt(28)
    r1.font.bold = False
    r1.font.color.rgb = _c(INK)
    r1.font.name = FONT_LIGHT
    if title_bold:
        r2 = p.add_run()
        r2.text = ("" if title_plain.endswith(" ") else " ") + title_bold
        r2.font.size = Pt(28)
        r2.font.bold = False        # Black уже даёт вес
        r2.font.color.rgb = _c(INK)
        r2.font.name = FONT_BLACK
    rect(slide, right - 1.6, 1.18, 1.6, 0.045, AMBER)
    if idx is not None:
        nb = slide.shapes.add_textbox(int(LEFT * EMU), int(0.4 * EMU),
                                      int(0.9 * EMU), int(0.6 * EMU))
        ntf = nb.text_frame
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        np = ntf.paragraphs[0]
        np.alignment = PP_ALIGN.LEFT
        np.text = f"{idx:02d}"
        np.font.size = Pt(20)
        np.font.color.rgb = _c(SOFT)
        np.font.name = FONT_LIGHT


def section_label(slide, left, top, w, label):
    tb = slide.shapes.add_textbox(int(left * EMU), int(top * EMU),
                                  int(w * EMU), int(0.3 * EMU))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p.text = label.upper()
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = _c(MUTED)
    _font(p.font)
    rect(slide, left, top + 0.32, w, 0.02, INK)


def top_right_brand(slide, text_str="Модель застройки территории"):
    """Бренд-плашка в правом верхнем углу (амбер-акцент)."""
    tb = slide.shapes.add_textbox(int(8.6 * EMU), int(0.22 * EMU),
                                  int(4.1 * EMU), int(0.3 * EMU))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p.text = text_str.upper()
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = _c(AMBER)
    _font(p.font)


def footer(slide, variant_name=None):
    """Нижний колонтитул — правовой дисклеймер (не проектная документация)."""
    rect(slide, LEFT, 7.12, CONTENT_W, 0.012, HAIR)
    tb = slide.shapes.add_textbox(int(LEFT * EMU), int(7.16 * EMU),
                                  int(CONTENT_W * EMU), int(0.28 * EMU))
    p = tb.text_frame.paragraphs[0]
    p.text = ("Информационная модель для оценки вариантов застройки. "
              "Не является проектной документацией и не влечёт правовых "
              "последствий. Значения — расчётные, требуют проверки проектом.")
    p.font.size = Pt(7.5)
    p.font.color.rgb = _c(SOFT)
    _font(p.font)


# Серые оттенки ушек: активное — тёмное и впритык к краю; прочие — светло-серые
# и чуть утоплены (как алфавитный указатель в записной книжке).
_TAB_ACTIVE = (0x4A, 0x4A, 0x4A)
_TAB_INACTIVE = (0xC4, 0xC4, 0xC4)
_RAIL_TOP = 1.55
_RAIL_BOT = 7.0
_EDGE = 13.28          # правая кромка листа (13.333 − поле)


def variant_rail(slide, labels: list[str], current: int | None) -> None:
    """Вертикальная «рейка» ушек всех вариантов с ФИКСИРОВАННЫМИ позициями:
    сверху вниз — База, Вариант 1, 2… Каждое ушко сохраняет свою высоту на всех
    слайдах. Активный (current) — тёмно-серый, впритык к правому краю и шире;
    прочие — светло-серые, чуть утоплены. Только оттенки серого."""
    n = len(labels)
    if n == 0:
        return
    span = _RAIL_BOT - _RAIL_TOP
    gap = 0.10
    slot = span / n
    tab_h = max(0.5, slot - gap)
    for i, label in enumerate(labels):
        active = (i == current)
        y = _RAIL_TOP + i * slot + (slot - tab_h) / 2
        w = 0.40 if active else 0.30
        x = _EDGE - w                      # впритык к правому краю
        color = _TAB_ACTIVE if active else _TAB_INACTIVE
        rect(slide, x, y, w, tab_h, color)
        cx, cy = x + w / 2, y + tab_h / 2
        half = tab_h / 2 + 0.6
        tb = slide.shapes.add_textbox(int((cx - half) * EMU), int((cy - 0.14) * EMU),
                                      int(2 * half * EMU), int(0.28 * EMU))
        tb.rotation = 270
        tf = tb.text_frame
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = label.upper()
        p.font.size = Pt(10 if active else 8.5)
        p.font.bold = active
        p.font.color.rgb = _c(WHITE)
        _font(p.font)


def kpi_card(slide, left, top, w, h, value, label, *, status=None,
             status_color=None, delta=None):
    rect(slide, left, top, w, h, WHITE, line=HAIR)
    tb = slide.shapes.add_textbox(int((left + 0.12) * EMU), int((top + 0.1) * EMU),
                                  int((w - 0.24) * EMU), int(0.5 * EMU))
    tf = tb.text_frame
    tf.word_wrap = False
    pv = tf.paragraphs[0]
    pv.text = str(value)
    pv.font.size = Pt(21)
    pv.font.bold = False
    pv.font.color.rgb = _c(INK)
    pv.font.name = FONT_LIGHT       # тонкие крупные цифры, как в модели
    lb = slide.shapes.add_textbox(int((left + 0.12) * EMU), int((top + 0.62) * EMU),
                                  int((w - 0.24) * EMU), int(0.3 * EMU))
    pl = lb.text_frame.paragraphs[0]
    pl.text = label.upper()
    pl.font.size = Pt(8.5)
    pl.font.color.rgb = _c(MUTED)
    _font(pl.font)
    extra = delta or status
    if extra:
        col = status_color or (OK if (delta and str(delta).startswith("+")) else MUTED)
        eb = slide.shapes.add_textbox(int((left + 0.12) * EMU),
                                      int((top + h - 0.32) * EMU),
                                      int((w - 0.24) * EMU), int(0.28 * EMU))
        pe = eb.text_frame.paragraphs[0]
        pe.text = str(extra)
        pe.font.size = Pt(9)
        pe.font.bold = True
        pe.font.color.rgb = _c(col)
        _font(pe.font)


def status_pill(slide, left, top, w, label, color):
    rect(slide, left, top, w, 0.32, WHITE, line=color, line_w=1.0)
    tb = slide.shapes.add_textbox(int(left * EMU), int((top + 0.01) * EMU),
                                  int(w * EMU), int(0.3 * EMU))
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = label
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = _c(color)
    _font(p.font)


def table(slide, left, top, w, headers, rows, *, col_ratios=None, fsize=11,
          row_colors=None, bold_last=False):
    nr, nc = len(rows) + 1, len(headers)
    h = 0.4 + 0.32 * len(rows)
    gt = slide.shapes.add_table(nr, nc, int(left * EMU), int(top * EMU),
                                int(w * EMU), int(h * EMU))
    tbl = gt.table
    tbl.first_row = True
    tbl.horz_banding = False
    if col_ratios:
        tot = sum(col_ratios)
        for j, rr in enumerate(col_ratios):
            tbl.columns[j].width = int(w * EMU * rr / tot)
    for j, htxt in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = str(htxt).upper()
        cell.fill.solid()
        cell.fill.fore_color.rgb = _c(INK)
        cell.margin_left = Emu(int(0.08 * EMU))
        cell.margin_top = Emu(int(0.02 * EMU))
        cell.margin_bottom = Emu(int(0.02 * EMU))
        for p in cell.text_frame.paragraphs:
            for r_ in p.runs:
                r_.font.size = Pt(fsize - 1)
                r_.font.bold = True
                r_.font.color.rgb = _c(WHITE)
                _font(r_.font)
    n = len(rows)
    for i, row in enumerate(rows, 1):
        txt_col = (row_colors[i - 1] if row_colors and i - 1 < len(row_colors)
                   else INK)
        is_last = bold_last and i == n
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _c(WHITE if i % 2 else ZEBRA)
            cell.margin_left = Emu(int(0.08 * EMU))
            cell.margin_top = Emu(int(0.015 * EMU))
            cell.margin_bottom = Emu(int(0.015 * EMU))
            for p in cell.text_frame.paragraphs:
                for r_ in p.runs:
                    r_.font.size = Pt(fsize)
                    r_.font.bold = is_last
                    r_.font.color.rgb = _c(txt_col)
                    _font(r_.font)
    return gt


def title_slide(deck, title_plain, title_bold, subtitle, meta_line):
    s = deck.slide()
    text(s, 0.9, 2.5, 11.5, 1.0, title_plain, size=40, font_name=FONT_LIGHT)
    text(s, 0.9, 3.35, 11.5, 1.0, title_bold, size=40, font_name=FONT_BLACK)
    rect(s, 0.92, 4.5, 2.2, 0.06, AMBER)
    text(s, 0.9, 4.75, 11.5, 0.6, subtitle, size=18, color=MUTED)
    text(s, 0.9, 6.6, 11.5, 0.5, meta_line, size=11, color=SOFT)
    return s
