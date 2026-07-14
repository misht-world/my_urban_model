"""Альбом КОНЦЕПЦИИ (PPTX) — мульти-вариантный (v0.13.0).

Структура (по ТЗ Михаила):
  1. Титул
  2. (Общая информация о территории — позже)
  3. Сводная сравнительная таблица (сначала — обзор всех вариантов)
  4. Базовый вариант: карточка (KPI как на сайте) + слайды-таблицы по вкладкам
  5. Варианты оптимизации: то же — карточка + те же таблицы

Таблицы берутся из ЕДИНОГО билдера `export/variant_tables.py` — те же строки,
что на «Расчёте», поэтому Базовый и оптимизационные варианты идентичны.
Экономика и полный аудит в альбом не входят (по запросу).
"""

from __future__ import annotations

import datetime as _dt

from urban_model.export.album import theme as T
from urban_model.export.table import KPI_SECTION_LABELS, results_to_dataframe
from urban_model.export.variant_tables import build_variant_table_blocks
from urban_model.models.result import TEPResult
from pptx.enum.text import PP_ALIGN

from urban_model.ui.formatting import fmt_int, fmt_m2

_MARGIN = T.LEFT
_CONTENT_W = T.CONTENT_W

# Значки KPI-карточек. Не эмодзи (те дают «тофу»-квадрат или цветной глиф), а
# ТЕКСТОВЫЕ символы Unicode из старых блоков (Геометрия/Дингбаты/Letterlike) —
# они гарантированно рендерятся МОНОХРОМНО и красятся цветом подписи (серым),
# как значки в модели.
_EMO = {
    "КИТ (ПЗЗ)": "⌂",          # ⌂ дом
    "Население, чел.": "☰",     # ☰ ряды
    "Площадь квартир": "▦",     # ▦ сетка
    "Этажность": "▤",           # ▤ слои
    "ЗНОП": "✿",               # ✿ флёрон
    "ДОО": "❍",                # ❍ круг
    "СОШ": "✎",                # ✎ карандаш
    "Доп. образование": "❖",    # ❖ ромб
    "Поликлиника": "✚",         # ✚ крест
    "Парковки": "◈",           # ◈ ромб-в-ромбе
    "Эконом-индекс": "₽",       # ₽ рубль
    "Выход жилья": "▲",         # ▲ вверх
    "Соц. нагрузка": "§",       # § параграф
    "ROI": "↗",                # ↗ рост
}


def _clean_name(name: str) -> str:
    """Убрать служебный префикс «opt:» из имени сценария."""
    return str(name).removeprefix("opt:").strip()


def _floors_txt(tep: TEPResult) -> str:
    ef = getattr(tep, "effective_floors", None) or 0
    if getattr(tep, "floor_clusters_detail", None):
        return f"{ef:.1f} ср."
    return f"{ef:.0f} эт." if ef else "—"


def _buckets_extra(formula: str | None, total: int) -> str | None:
    """«N объектов по lo–hi мест» — та же подпись, что под метриками ДОО/СОШ
    в KPI-сетке модели (ui/output._buckets_delta)."""
    from urban_model.export.variant_tables import parse_object_buckets
    if not formula or total == 0:
        return None
    n, caps = parse_object_buckets(formula)
    if n == 0:
        return None
    lo, hi = min(caps), max(caps)
    per = f"{lo} мест" if lo == hi else f"{lo}–{hi} мест"
    word = "объект" if n == 1 else ("объекта" if 2 <= n <= 4 else "объектов")
    return f"{n} {word} по {per}"


def _kpi_pairs(tep: TEPResult) -> list[tuple[str, str, str | None]]:
    """(значение, подпись с эмодзи, доп.строка) — зеркало render_main_kpi_grid."""
    def lbl(t):
        e = _EMO.get(t)
        return f"{e}  {t}" if e else t

    pairs: list[tuple[str, str, str | None]] = []
    pairs.append((f"{tep.kit.value:.3f}", lbl("КИТ (ПЗЗ)"), None))
    pairs.append((f"{fmt_int(tep.population.value)}", lbl("Население, чел."), None))
    pairs.append((fmt_m2(tep.apartments_area.value), lbl("Площадь квартир"), None))
    pairs.append((_floors_txt(tep), lbl("Этажность"), None))
    znop_pp = tep.znop_per_person.value or 0
    znop_area = int(tep.znop_area.value or 0)
    pairs.append((f"{znop_area:,}".replace(",", " ") + " м²" if znop_area else "—",
                  lbl("ЗНОП"), f"{znop_pp:.1f} м²/чел" if znop_pp > 0 else None))
    kg = int(tep.kindergarten_places_accepted.value or 0)
    pairs.append((f"{kg} мест" if kg else "—", lbl("ДОО"),
                  _buckets_extra(tep.kindergarten_places_accepted.formula, kg)))
    sch = int(tep.school_places_accepted.value or 0)
    pairs.append((f"{sch} мест" if sch else "—", lbl("СОШ"),
                  _buckets_extra(tep.school_places_accepted.formula, sch)))
    ae = int(tep.add_education_places_accepted.value or 0)
    ae_place = ("встроен. (ВПП)" if getattr(tep, "add_education_built_in", False)
                else "отд. стоящее")
    pairs.append((f"{ae} мест" if ae else "—", lbl("Доп. образование"),
                  ae_place if ae else None))
    poly_f = getattr(tep, "polyclinic_visits_accepted", None)
    poly = int(poly_f.value or 0) if poly_f is not None else 0
    poly_place = ("ВПП (офис врача)" if getattr(tep, "polyclinic_built_in", False)
                  else "отд. стоящая")
    pairs.append((f"{poly} посещ." if poly else "—", lbl("Поликлиника"),
                  poly_place if poly else None))
    op = int(tep.parking_open_places.value or 0)
    ml = int(tep.parking_multilevel_places.value or 0)
    ug = int(tep.parking_underground_places.value or 0)
    st_f = getattr(tep, "parking_stylobate_places", None)
    styl = int(st_f.value or 0) if st_f is not None else 0
    total_pl = int(tep.parking_required_places.value or 0)
    parts = []
    if op:   parts.append(f"откр. {op}")
    if ml:   parts.append(f"мн. {ml}")
    if styl: parts.append(f"ст. {styl}")
    if ug:   parts.append(f"пд. {ug}")
    pairs.append((f"{total_pl} м/м" if total_pl else "—", lbl("Парковки"),
                  " · ".join(parts) or None))
    return pairs


def _econ_pairs(tep: TEPResult) -> list[tuple[str, str, str | None]]:
    e = getattr(tep, "economy", None)
    if e is None:
        return []

    def lbl(t):
        e = _EMO.get(t)
        return f"{e}  {t}" if e else t

    out: list[tuple[str, str, str | None]] = [
        (f"{e.economy_index:.0f} / 100", lbl("Эконом-индекс"), None),
        (f"{e.sellable_ratio * 100:.0f}%" if e.sellable_ratio else "—",
         lbl("Выход жилья"), None),
    ]
    social_cost = (e.cost.kindergarten + e.cost.school + e.cost.social_parking
                   + getattr(e.cost, "add_education", 0.0))
    if social_cost > 0.5 or e.revenue.social_compensation > 0.5:
        out.append((f"{-e.net_social_burden:+,.0f}".replace(",", " "),
                    lbl("Соц. нагрузка"), "баллы"))
    else:
        out.append((f"{e.roi * 100:.1f}%" if e.cost.total > 0 else "—", lbl("ROI"), None))
    return out


def _phasing_chip(tep: TEPResult) -> str | None:
    """«Очерёдность: 3 очереди (30/40/30%)» — если очереди заданы."""
    ph = getattr(tep, "phasing", None)
    if ph is None or not ph.stages:
        return None
    shares = "/".join(f"{s.share * 100:.0f}" for s in ph.stages)
    n = len(ph.stages)
    word = "очереди" if 2 <= n <= 4 else "очередей"
    n_lots = max((s.lot for s in ph.stages), default=1)
    lots = f" · {n_lots} лот(а)" if n_lots > 1 else ""
    ok = all(s.is_ok for s in ph.stages)
    status = "" if ok else " · есть дефициты"
    auto = " · авто" if getattr(ph, "mode", "") == "auto" else ""
    return f"Очерёдность: {n} {word} ({shares}%){lots}{auto}{status}"


def _var_label(v_index: int) -> str:
    return "База" if v_index == 0 else f"Вариант {v_index}"


# ── Слайд «Общая информация о территории» (из CalculationOptions Базы) ──

_PARK_MODE_RU = {
    "min_open": "минимум открытых (12.5%), остальное — подземные",
    "all_open": "100% открытые",
    "custom": "заданные доли",
}
_FUND_RU = {
    "compensated": "компенсация города (доля)",
    "developer": "за счёт застройщика",
    "city": "строит город",
    "at_cost": "передача по себестоимости",
}


def _territory_rows(site_area: float, o) -> list[tuple[str, str]]:
    """Строки «Исходные условия» из options (левая таблица)."""
    rows = [("Площадь квартала", f"{fmt_m2(site_area)}  ({site_area/10_000:.2f} га)")]
    if getattr(o, "floor_clusters", None):
        zones = " + ".join(f"{c.area_m2/10_000:.1f} га × {c.floors} эт."
                           for c in o.floor_clusters)
        rows.append(("Этажность (зоны)", zones))
    else:
        rows.append(("Этажность", f"{o.floors} эт."))
    rows.append(("Документация по планировке (ДПТ)",
                 "да (КИТ ≤ 2.5)" if o.planning_doc else "нет (КИТ ≤ 1.4)"))
    rows.append(("Класс жилья", {"economy": "эконом", "comfort": "комфорт",
                                 "business": "бизнес"}.get(o.residential_class,
                                                           str(o.residential_class))))
    p = o.parking
    park = _PARK_MODE_RU.get(p.mode, p.mode)
    if p.mode == "custom":
        parts = []
        if p.open_share > 0.001:
            parts.append(f"откр. {p.open_share:.0%}")
        if p.multilevel_share > 0.001:
            parts.append(f"многоур. {p.multilevel_share:.0%} ({p.multilevel_levels} эт.)")
        if p.underground_share > 0.001:
            parts.append(f"подз. {p.underground_share:.0%} ({p.underground_levels} ур.)")
        if getattr(p, "stylobate_share", 0) > 0.001:
            parts.append(f"стилобат {p.stylobate_share:.0%}")
        park = " · ".join(parts)
    rows.append(("Парковки (жилищные)", park))
    vpp_total = sum(b.area_m2 for b in (o.built_in_list or []))
    if o.built_in is not None:
        vpp_total += o.built_in.area_m2
    if vpp_total > 0:
        rows.append(("ВПП (встроенно-пристроенные)", fmt_m2(vpp_total)))
    if o.znop_per_person_override is not None:
        rows.append(("ЗНОП (вручную)", f"{o.znop_per_person_override:g} м²/чел"))
    if o.znop_total_area_override is not None:
        rows.append(("ЗНОП, площадь (вручную)", fmt_m2(o.znop_total_area_override)))
    # v0.15.9 (п.7): остальные данные левых карточек — проезды и очерёдность.
    if getattr(o, "driveways_intra_share_override", None) is not None:
        rows.append(("Внутриквартальные проезды (вручную)",
                     f"{o.driveways_intra_share_override:.0%} от квартала"))
    if getattr(o, "driveways_lot_share_override", None) is not None:
        rows.append(("Проезды на ЗУ жилья (вручную)",
                     f"{o.driveways_lot_share_override:.0%}"))
    _ph = getattr(o, "phasing", None)
    if _ph is not None:
        _mode = ("авто — по обеспеченности соцобъектами" if _ph.mode == "auto"
                 else f"{len(_ph.shares)} очереди вручную ("
                      + "/".join(f"{s * 100:.0f}" for s in _ph.shares) + "%)")
        if getattr(_ph, "engineering_by_lots", False):
            _mode += " · инженерия автономно по лотам"
        rows.append(("Очерёдность застройки", _mode))
    if getattr(o, "include_economy", True):
        rows.append(("Соцобъекты (финансирование)",
                     _FUND_RU.get(getattr(o, "social_funding", "compensated"), "—")))
    return rows


def _composition_rows(o) -> list[tuple[str, str]]:
    """Строки «Состав расчёта» — что учитывается (правая таблица)."""
    def yn(v):
        return "учитывается" if v else "не учитывается"

    def od(spec):
        return " (только потребность)" if getattr(spec, "only_demand", False) else ""

    rows = [
        ("ДОО (детские сады)", yn(o.include_kindergarten) + od(o.kindergarten)
         + (" · только типовые вместимости"
            if getattr(o.kindergarten, "strict_capacity", False) else "")),
        ("СОШ (школы)", yn(o.include_school) + od(o.school)
         + (" · только типовые вместимости"
            if getattr(o.school, "strict_capacity", False) else "")),
        ("Доп. образование (ВРИ 3.5.1)", yn(getattr(o, "include_add_education", True))),
        ("Поликлиника (ВРИ 3.4.1)", yn(getattr(o, "include_polyclinic", True))),
        ("Спортивные сооружения (ВРИ 5.1.3)", yn(o.include_sport_facilities)),
        ("ЗНОП", yn(o.include_znop)),
        ("Парковки", yn(o.include_parking)),
        ("Внутриквартальные проезды", yn(o.include_intra_driveways)),
        ("Инженерная инфраструктура", yn(getattr(o, "include_engineering", True))),
        ("Экономика (усл. баллы)", yn(getattr(o, "include_economy", True))),
        ("Норматив озеленения 25%",
         "контролируется" if o.enforce_quarter_greening_norm else "справочно"),
        ("Норматив плотности 450 чел/га",
         "контролируется" if o.enforce_density_norm else "справочно"),
    ]
    if getattr(o, "custom_objects", None):
        rows.append(("Пользовательские объекты", f"{len(o.custom_objects)} шт."))
    return rows


def _territory_slide(deck, site_area: float, o, rail_labels: list[str]) -> None:
    """Слайд «Общая информация о территории»: исходные условия + состав расчёта."""
    s = deck.slide()
    T.title_band(s, "Общая информация о", "территории")
    T.top_right_brand(s)
    half = (_CONTENT_W - 0.3) / 2
    T.section_label(s, _MARGIN, 1.5, half, "Исходные условия")
    T.table(s, _MARGIN, 1.95, half, ["Параметр", "Значение"],
            _territory_rows(site_area, o), col_ratios=[1.1, 1.4], fsize=10)
    right_x = _MARGIN + half + 0.3
    T.section_label(s, right_x, 1.5, half, "Состав расчёта")
    T.table(s, right_x, 1.95, half, ["Компонент", "Статус"],
            _composition_rows(o), col_ratios=[1.4, 1.1], fsize=10)
    T.footer(s)


def _split_title(title: str) -> tuple[str, str]:
    """Разбить заголовок на (обычная часть, жирное последнее слово)."""
    parts = title.rsplit(" ", 1)
    if len(parts) == 2:
        return parts[0], parts[1]
    return "", title


def _chrome(s, rail_labels: list[str], current: int | None,
            accents: dict | None = None) -> None:
    """Общие элементы слайда: бренд справа-сверху + рейка вариантов справа."""
    T.top_right_brand(s)
    T.variant_rail(s, rail_labels, current, accents=accents)


def _variant_accents(scenarios: list[tuple[str, TEPResult]]) -> dict:
    """Слабые цветовые акценты вариантов (v0.17.2, п.8): База — синий,
    «Девелоперский» (по имени сценария) — амбер."""
    acc: dict = {0: T.ACCENT_BASE}
    for i, (nm, _) in enumerate(scenarios):
        if "девелоперск" in str(nm).lower():
            acc[i] = T.ACCENT_DEV
    return acc


def _card_slide(deck, name: str, tep: TEPResult, kind: str, v_index: int,
                rail_labels: list[str], accents: dict | None = None) -> None:
    """Слайд-карточка варианта: KPI-сетка как на сайте (2 ряда × 5 + экономика)."""
    s = deck.slide()
    # Заголовок с частичным жирным (как «Сравнение вариантов»): последнее слово
    # — Black. «Базовый вариант» → Базовый + **вариант**; «Вариант 1» → Вариант + **1**.
    _plain, _bold = _split_title(kind)
    T.title_band(s, _plain, _bold)
    _chrome(s, rail_labels, v_index, accents)
    # Длинное имя — отдельной строкой под чертой (перенос), по правому краю.
    # Единый стиль имени варианта на всех карточках (как на слайдах-таблицах).
    T.text(s, _MARGIN, 1.3, _CONTENT_W, 0.6, _clean_name(name),
           size=11, color=T.MUTED, spacing=1.05, align=PP_ALIGN.RIGHT)

    def _draw_row(pairs, top):
        n = len(pairs)
        if n == 0:
            return
        gap = 0.12
        cw = (_CONTENT_W - gap * (n - 1)) / n
        for i, (val, lbl, extra) in enumerate(pairs):
            left = _MARGIN + i * (cw + gap)
            # v0.17.2 (п.9): «Площадь квартир» — слабый амбер-«текстовыделитель»
            # за значением.
            _hl = "Площадь квартир" in lbl
            T.kpi_card(s, left, top, cw, 1.15, val, lbl, delta=extra,
                       value_hl=_hl)

    kp = _kpi_pairs(tep)
    _draw_row(kp[:5], 1.95)
    _draw_row(kp[5:10], 3.25)
    econ = _econ_pairs(tep)
    if econ:
        T.section_label(s, _MARGIN, 4.6, _CONTENT_W, "Экономика")
        _draw_row(econ, 5.0)
    # Очерёдность (v0.15.1; v0.16.2 п.9 — перенесена ВНИЗ карточки).
    _chip = _phasing_chip(tep)
    if _chip:
        T.text(s, _MARGIN, 6.6, _CONTENT_W, 0.35, _chip, size=10, color=T.MUTED)
    T.footer(s)


def _table_slide(deck, name: str, block, v_index: int,
                 rail_labels: list[str], accents: dict | None = None) -> None:
    """Слайд-таблица одной вкладки варианта (из общего билдера) + итог раздела."""
    s = deck.slide()
    T.title_band(s, block.title, "")
    _chrome(s, rail_labels, v_index, accents)
    T.text(s, _MARGIN, 1.28, _CONTENT_W, 0.35, _clean_name(name),
           size=11, color=T.MUTED, align=PP_ALIGN.RIGHT)
    if block.columns is None:
        headers = ["Показатель", "Значение"]
        rows = [(r["Показатель"], r["Значение"]) for r in block.rows]
        ratios = [3, 2]
    else:
        # v0.15.1: для альбома — компактный набор колонок (album_columns),
        # если задан: широкие таблицы (очерёдность, 10 колонок) на слайде
        # переносят текст в ячейках, реальная высота строк растёт и подписи
        # налезали на таблицу.
        cols = getattr(block, "album_columns", None) or block.columns
        headers = cols
        rows = [[r.get(c, "") for c in cols] for r in block.rows]
        ratios = None
    fsize = 11 if len(headers) <= 6 else 9.5
    row_h = 0.32 if len(headers) <= 6 else 0.48
    # v0.15.10 (п.2 Михаила): длинные таблицы (баланс с зонами и т.п.)
    # ПЕРЕНОСЯТСЯ на следующие слайды, а не наезжают на подписи/подвал.
    page_rows = 13 if len(headers) <= 6 else 9
    top = 1.75
    pages = [rows[i:i + page_rows] for i in range(0, len(rows), page_rows)] or [[]]
    for pi, page in enumerate(pages):
        if pi > 0:
            s = deck.slide()
            T.title_band(s, block.title, "(продолжение)")
            _chrome(s, rail_labels, v_index, accents)
            T.text(s, _MARGIN, 1.28, _CONTENT_W, 0.35, _clean_name(name),
                   size=11, color=T.MUTED, align=PP_ALIGN.RIGHT)
        T.table(s, _MARGIN, top, _CONTENT_W, headers, page,
                col_ratios=ratios, fsize=fsize)
        note_top = top + 0.4 + row_h * len(page) + 0.12
        if pi == len(pages) - 1:
            # Итог и примечания — только на последней странице таблицы.
            if block.summary:
                T.text(s, _MARGIN, note_top, _CONTENT_W, 0.6, block.summary,
                       size=9, color=T.MUTED)
                note_top += 0.42
            for note in block.notes:
                T.text(s, _MARGIN, note_top, _CONTENT_W, 0.6, note,
                       size=9, color=T.MUTED)
                note_top += 0.42
        T.footer(s)


def _charts_slide(deck, name: str, tep: TEPResult, v_index: int,
                  rail_labels: list[str], accents: dict | None = None) -> None:
    """Слайд «Диаграммы» в конце варианта (v0.17.2, п.10 Михаила):
    баланс территории (стек-полоса) + удельные показатели на жителя."""
    from urban_model.export.album import charts as C
    buf_bal = C.chart_balance(tep)
    buf_pc = C.chart_per_capita(tep)
    if buf_bal is None and buf_pc is None:
        return
    s = deck.slide()
    T.title_band(s, "Диаграммы", "")
    _chrome(s, rail_labels, v_index, accents)
    T.text(s, _MARGIN, 1.28, _CONTENT_W, 0.35, _clean_name(name),
           size=11, color=T.MUTED, align=PP_ALIGN.RIGHT)
    # Две диаграммы рядом: слева широкая полоса баланса, справа удельные.
    y = 1.95
    if buf_bal is not None:
        T.section_label(s, _MARGIN, y, 6.9, "Баланс территории")
        s.shapes.add_picture(buf_bal, int(_MARGIN * T.EMU),
                             int((y + 0.4) * T.EMU),
                             width=int(6.9 * T.EMU))
    if buf_pc is not None:
        _x = _MARGIN + 7.25
        T.section_label(s, _x, y, _CONTENT_W - 7.25, "На одного жителя")
        s.shapes.add_picture(buf_pc, int(_x * T.EMU),
                             int((y + 0.4) * T.EMU),
                             width=int(4.3 * T.EMU))
    T.footer(s)


_WC_TAG_RE = None


def _clean_warning(w: str) -> str:
    """Убрать служебный тег кода «[CODE] …» из текста предупреждения."""
    import re
    global _WC_TAG_RE
    if _WC_TAG_RE is None:
        _WC_TAG_RE = re.compile(r"^\[[A-Z0-9_]+\]\s*")
    return _WC_TAG_RE.sub("", str(w)).strip()


def _notes_slide(deck, name: str, tep: TEPResult, v_index: int,
                 rail_labels: list[str], accents: dict | None = None) -> None:
    """Лист «Примечания» в конце варианта (v0.17.2, п.5 Михаила):
    предупреждения и информационные пометки расчёта (те же, что показываются
    на «Расчёте»). Если предупреждений нет — лист не создаётся."""
    notes = [_clean_warning(w) for w in (tep.warnings or []) if str(w).strip()]
    if not notes:
        return
    per_page = 9
    pages = [notes[i:i + per_page] for i in range(0, len(notes), per_page)]
    for pi, page in enumerate(pages):
        s = deck.slide()
        T.title_band(s, "Примечания", "(продолжение)" if pi else "")
        _chrome(s, rail_labels, v_index, accents)
        T.text(s, _MARGIN, 1.28, _CONTENT_W, 0.35, _clean_name(name),
               size=11, color=T.MUTED, align=PP_ALIGN.RIGHT)
        y = 1.85
        for w in page:
            T.text(s, _MARGIN, y, _CONTENT_W, 0.55, "•  " + w,
                   size=10, color=T.INK, spacing=1.05)
            y += 0.56
        T.footer(s)


def _divider_slide(deck, plain: str, bold: str) -> None:
    """Лист-прокладка (раздел альбома): крупный заголовок по центру листа,
    амбер-черта — в стиле титульного (v0.16.2, п.1 Михаила)."""
    s = deck.slide()
    T.top_right_brand(s)
    T.text(s, 0.9, 2.85, 11.5, 0.85, plain, size=34, font_name=T.FONT_LIGHT)
    T.text(s, 0.9, 3.6, 11.5, 0.85, bold, size=34, font_name=T.FONT_BLACK)
    T.rect(s, 0.92, 4.62, 2.2, 0.06, T.AMBER)
    T.footer(s)


def _reorder_blocks_for_album(blocks: list) -> list:
    """Порядок листов в альбоме: «Баланс территории» — ПЕРВЫМ, перед «Жильём»
    (v0.17.2, п.7 Михаила; в v0.16.2 был после «Жилья»). На «Расчёте» порядок
    вкладок остаётся прежним (правка только для альбома)."""
    bal = next((b for b in blocks if getattr(b, "key", "") == "balance"), None)
    if bal is None:
        return blocks
    return [bal] + [b for b in blocks if b is not bal]


# Показатели, где БОЛЬШЕ = лучше (подсвечиваем максимум в строке).
_HL_MORE_IS_BETTER = ("Площадь квартир", "Эконом-индекс")


def _best_col_for(df, lab: str, names: list[str]) -> int | None:
    """Индекс варианта (0-based) с максимальным числовым значением в строке."""
    best_j, best_v = None, None
    for j, nm in enumerate(names):
        try:
            v = float(str(df.loc[lab, nm]).replace(" ", "").replace(" ", ""))
        except (ValueError, TypeError):
            continue
        if best_v is None or v > best_v:
            best_v, best_j = v, j
    return best_j


def _comparison_slides(deck, scenarios: list[tuple[str, TEPResult]],
                       rail_labels: list[str],
                       accents: dict | None = None) -> None:
    """Сводная сравнительная таблица (как в xlsx). Разбита по слайдам."""
    clean = [(_clean_name(n), t) for n, t in scenarios]
    df = results_to_dataframe(clean)
    names = list(df.columns)
    labels = list(df.index)
    # Лучший вариант (столбец) по каждому «больше=лучше» показателю.
    best_by_label = {
        lab: _best_col_for(df, lab, names)
        for lab in labels if any(k in lab for k in _HL_MORE_IS_BETTER)
    }
    ratios = [2.4] + [1.0] * len(names)
    top = 1.7
    chunk = 13
    for start in range(0, len(labels), chunk):
        s = deck.slide()
        part = "" if len(labels) <= chunk else f" ({start // chunk + 1})"
        T.title_band(s, "Сравнение", "вариантов" + part)
        _chrome(s, rail_labels, None, accents)   # рейка со всеми вариантами
        sub = labels[start:start + chunk]
        headers = ["Показатель"] + names
        rows = []
        hl_cells = set()
        section_rows = set()
        for i, lab in enumerate(sub, 1):     # i — 1-based строка данных в таблице
            row = [lab]
            for nm in names:
                v = df.loc[lab, nm]
                row.append("" if v is None else str(v))
            rows.append(row)
            # v0.16.2 (п.4): строки-заголовки категорий (ДОО/СОШ/ЗНОП/парковки…)
            # — серая заливка + жирный (см. theme.table section_rows).
            if lab in KPI_SECTION_LABELS:
                section_rows.add(i)
                continue
            bj = best_by_label.get(lab)
            if bj is not None:
                hl_cells.add((i, bj + 1))    # +1: столбец 0 — «Показатель»
        # (item 3) Серые подписи «База»/«Вариант N» над колонками — для сопоставления
        # с боковой рейкой. Позиции колонок из ratios.
        _tot = sum(ratios)
        for j in range(len(names)):
            x0 = _MARGIN + _CONTENT_W * (2.4 + j) / _tot
            cw = _CONTENT_W * 1.0 / _tot
            T.text(s, x0, top - 0.28, cw, 0.24, _var_label(j),
                   size=8, color=T.SOFT, align=PP_ALIGN.CENTER)
            # v0.17.2 (п.8): слабая цветовая полоска под подписью варианта
            # (синий — База, амбер — Девелоперский).
            _acc = (accents or {}).get(j)
            if _acc is not None:
                T.rect(s, x0 + cw * 0.2, top - 0.055, cw * 0.6, 0.04, _acc)
        T.table(s, _MARGIN, top, _CONTENT_W, headers, rows, col_ratios=ratios,
                fsize=10, hl_cells=hl_cells, section_rows=section_rows)
        T.footer(s)


def build_concept_album(
    scenarios: list[tuple[str, TEPResult]], path: str, *,
    title_line: str = "Альбом",
    title_bold: str = "концепции",
    base_options=None,
    site_area: float | None = None,
) -> str:
    """Собрать мульти-вариантный альбом концепции (PPTX 16:9). Возвращает путь.

    scenarios: список (имя, TEPResult) — Базовый вариант первым, далее варианты.
    base_options/site_area: параметры Базы для слайда «Общая информация о
    территории» (если не переданы — слайд пропускается).

    Структура (v0.13.5): Титул → Общая информация → Сравнение →
    карточки ВСЕХ вариантов подряд → детальные таблицы по каждому варианту.
    Каждый слайд обёрнут в try/except — альбом не падает на нехватке данных.
    """
    deck = T.Deck()
    n_var = len(scenarios)
    meta = (f"Вариантов: {n_var}   ·   {_dt.date.today().strftime('%d.%m.%Y')}"
            f"   ·   профиль СПб")
    T.title_slide(deck, title_line, title_bold,
                  subtitle="Базовый вариант и варианты оптимизации", meta_line=meta)

    # Фиксированные подписи ушек-рейки: сверху вниз — База, Вариант 1, 2…
    rail_labels = [_var_label(i) for i in range(n_var)]
    # Слабые цветовые акценты (v0.17.2, п.8): База — синий, Девелоперский — амбер.
    accents = _variant_accents(scenarios)

    # Общая информация о территории (из параметров Базы).
    if base_options is not None and site_area:
        try:
            _territory_slide(deck, float(site_area), base_options, rail_labels)
        except Exception:  # noqa: BLE001 — §12: альбом не должен падать
            pass

    # Сводное сравнение (обзор всех вариантов).
    try:
        _comparison_slides(deck, scenarios, rail_labels, accents)
    except Exception:  # noqa: BLE001
        pass

    # Карточки ВСЕХ вариантов подряд (обзор каждого на одном слайде).
    for v_idx, (name, tep) in enumerate(scenarios):
        kind = "Базовый вариант" if v_idx == 0 else f"Вариант {v_idx}"
        try:
            _card_slide(deck, name, tep, kind, v_idx, rail_labels, accents)
        except Exception:  # noqa: BLE001
            pass

    # Лист-прокладка перед детальными таблицами (v0.16.2, п.1).
    try:
        _divider_slide(deck, "Аналитическая", "часть")
    except Exception:  # noqa: BLE001
        pass

    # Детальные таблицы — по каждому варианту; в конце варианта — диаграммы
    # (баланс + удельные на жителя) и лист примечаний (v0.17.2, пп.5/10).
    for v_idx, (name, tep) in enumerate(scenarios):
        try:
            blocks = _reorder_blocks_for_album(build_variant_table_blocks(tep))
        except Exception:  # noqa: BLE001
            blocks = []
        for blk in blocks:
            try:
                _table_slide(deck, name, blk, v_idx, rail_labels, accents)
            except Exception:  # noqa: BLE001
                pass
        try:
            _charts_slide(deck, name, tep, v_idx, rail_labels, accents)
        except Exception:  # noqa: BLE001
            pass
        try:
            _notes_slide(deck, name, tep, v_idx, rail_labels, accents)
        except Exception:  # noqa: BLE001
            pass

    return deck.save(path)
