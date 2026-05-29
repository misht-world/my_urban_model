"""PPTX-презентация по варианту застройки (v0.10.10).

`build_variant_pptx(name, tep, path)` — слайды 16:9 для Заказчика:
  1. Титул;
  2. Сводка — все ключевые показатели кратко (ТЭП, парковки по типам,
     ДОО/СОШ, ЗНОП, ОКС, кластеры);
  3. Баланс территории (диаграмма + ключевое);
  4. Детально: жильё и социальная инфраструктура (таблица);
  5. Детально: парковки, проезды, кластеры (таблицы);
  6. Экономика (простая диаграмма себестоимость→выручка→прибыль + метрики).

Без теней. Переиспользует matplotlib-диаграммы из docx_report.
"""

from __future__ import annotations

import datetime as _dt
import re as _re

from urban_model.export.docx_report import (
    _chart_balance,
    _chart_economy,
    _fmt,
    _general_kpis,
)
from urban_model.models.result import TEPResult

_BLUE = (0x1F, 0x4E, 0x79)
_BLUE_LT = (0x2E, 0x75, 0xB6)
_GREY = (0x59, 0x59, 0x59)
_LIGHT = (0xEA, 0xF1, 0xF8)
_ZEBRA = (0xF2, 0xF6, 0xFB)
_WHITE = (0xFF, 0xFF, 0xFF)
EMU = 914400


def _buckets(formula: str | None) -> list[int]:
    if not formula:
        return []
    m = _re.search(r"\[([^\]]+)\]", formula)
    if not m:
        return []
    try:
        return [int(x.strip()) for x in m.group(1).split(",") if x.strip()]
    except ValueError:
        return []


def _objs_text(formula: str | None, total: int, unit: str = "мест") -> str:
    b = _buckets(formula)
    if not b or total == 0:
        return f"{total} {unit}"
    lo, hi = min(b), max(b)
    cap = f"{lo}" if lo == hi else f"{lo}–{hi}"
    return f"{total} {unit} ({len(b)} об. × {cap})"


def build_variant_pptx(name: str, tep: TEPResult, path: str) -> str:
    """Слайдовый отчёт (16:9) по одному варианту. Возвращает путь."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Emu, Pt

    SW, SH = int(13.333 * EMU), int(7.5 * EMU)
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    blank = prs.slide_layouts[6]

    def C(t):
        return RGBColor(*t)

    def _no_shadow(shape):
        try:
            shape.shadow.inherit = False
        except Exception:  # noqa: BLE001
            pass

    def _slide():
        return prs.slides.add_slide(blank)

    def _rect(slide, left, top, w, h, fill, line=None):
        sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    int(left * EMU), int(top * EMU),
                                    int(w * EMU), int(h * EMU))
        sp.fill.solid(); sp.fill.fore_color.rgb = C(fill)
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = C(line); sp.line.width = Pt(0.75)
        _no_shadow(sp)
        return sp

    def _band(slide, title, idx=None):
        bar = _rect(slide, 0, 0, 13.333, 0.95, _BLUE)
        # квадрат-бейдж слева + амбер-акцент под шапкой
        _rect(slide, 0.42, 0.3, 0.35, 0.35, (0xF5, 0xA6, 0x23))
        _rect(slide, 0, 0.95, 13.333, 0.05, (0xF5, 0xA6, 0x23))
        tb = slide.shapes.add_textbox(int(1.0 * EMU), int(0.12 * EMU),
                                      int(11.5 * EMU), int(0.72 * EMU))
        tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.text = title
        p.font.size = Pt(23); p.font.bold = True; p.font.color.rgb = C(_WHITE)
        if idx is not None:
            nb = slide.shapes.add_textbox(int(12.2 * EMU), int(0.18 * EMU),
                                          int(0.9 * EMU), int(0.6 * EMU))
            ntf = nb.text_frame; ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
            np = ntf.paragraphs[0]; np.alignment = PP_ALIGN.RIGHT
            np.text = f"{idx:02d}"
            np.font.size = Pt(22); np.font.bold = True
            np.font.color.rgb = C((0x9D, 0xC3, 0xE6))

    def _footer(slide, name):
        _rect(slide, 0, 7.18, 13.333, 0.02, (0xD6, 0xE4, 0xF0))
        fb = slide.shapes.add_textbox(int(0.42 * EMU), int(7.2 * EMU),
                                      int(12.5 * EMU), int(0.28 * EMU))
        p = fb.text_frame.paragraphs[0]
        p.text = f"Модель застройки территории  ·  вариант «{name}»"
        p.font.size = Pt(8); p.font.color.rgb = C((0x9A, 0xA7, 0xB4))

    def _text(slide, left, top, w, h, text, size=13, bold=False,
              color=_GREY, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(int(left * EMU), int(top * EMU),
                                      int(w * EMU), int(h * EMU))
        tf = tb.text_frame; tf.word_wrap = True
        lines = text.split("\n")
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = ln; p.alignment = align
            p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = C(color)
        return tb

    def _kpi(slide, left, top, w, h, value, label, accent=(0x2E, 0x75, 0xB6)):
        card = _rect(slide, left, top, w, h, _WHITE, line=(0xD4, 0xDF, 0xEC))
        _rect(slide, left, top, w, 0.07, accent)  # верхняя акцентная полоса
        tf = card.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pv = tf.paragraphs[0]; pv.alignment = PP_ALIGN.CENTER; pv.text = value
        pv.font.size = Pt(20); pv.font.bold = True; pv.font.color.rgb = C(_BLUE)
        pl = tf.add_paragraph(); pl.alignment = PP_ALIGN.CENTER; pl.text = label
        pl.font.size = Pt(10); pl.font.color.rgb = C(_GREY)

    def _table(slide, left, top, w, headers, rows, col_ratios=None, fsize=11):
        nr, nc = len(rows) + 1, len(headers)
        h = (0.45 + 0.34 * len(rows))
        gt = slide.shapes.add_table(nr, nc, int(left * EMU), int(top * EMU),
                                    int(w * EMU), int(h * EMU))
        tbl = gt.table
        tbl.first_row = True; tbl.horz_banding = False
        if col_ratios:
            tot = sum(col_ratios)
            for j, rr in enumerate(col_ratios):
                tbl.columns[j].width = int(w * EMU * rr / tot)
        # header
        for j, htxt in enumerate(headers):
            cell = tbl.cell(0, j); cell.text = str(htxt)
            cell.fill.solid(); cell.fill.fore_color.rgb = C(_BLUE)
            cell.margin_left = Emu(int(0.08 * EMU)); cell.margin_top = Emu(int(0.02 * EMU))
            for p in cell.text_frame.paragraphs:
                for r_ in p.runs:
                    r_.font.size = Pt(fsize); r_.font.bold = True; r_.font.color.rgb = C(_WHITE)
        # rows
        for i, row in enumerate(rows, 1):
            for j, val in enumerate(row):
                cell = tbl.cell(i, j); cell.text = str(val)
                cell.fill.solid()
                cell.fill.fore_color.rgb = C(_ZEBRA if i % 2 else _WHITE)
                cell.margin_left = Emu(int(0.08 * EMU)); cell.margin_top = Emu(int(0.02 * EMU))
                for p in cell.text_frame.paragraphs:
                    for r_ in p.runs:
                        r_.font.size = Pt(fsize); r_.font.color.rgb = C((0x22, 0x22, 0x22))
        return gt

    site_area = tep.balance.site_area

    # ─── Слайд 1: Титул ───
    s = _slide()
    _rect(s, 0, 0, 13.333, 7.5, _BLUE)
    _rect(s, 0, 5.0, 13.333, 0.06, (0xF5, 0xA6, 0x23))  # акцентная линия
    _text(s, 0.9, 2.3, 11.5, 1.4,
          "Технико-экономические показатели застройки",
          size=34, bold=True, color=_WHITE)
    _text(s, 0.9, 3.9, 11.5, 0.7, f"Вариант: «{name}»", size=20, color=(0xD6, 0xE4, 0xF0))
    _text(s, 0.9, 5.2, 11.5, 0.6,
          f"Площадь квартала: {_fmt(site_area)} м² ({site_area/10000:.2f} га)    ·    "
          f"{_dt.date.today().strftime('%d.%m.%Y')}    ·    профиль СПб",
          size=13, color=(0xBD, 0xD7, 0xEE))

    # ─── Слайд 2: Сводка ───
    s = _slide(); _band(s, "Сводка по варианту", idx=1); _footer(s, name)
    kpis = _general_kpis(tep)
    cw, ch, gx = 3.0, 1.2, 0.13
    x0, y0 = 0.5, 1.25
    for i, (label, value) in enumerate(kpis[:8]):
        r, c = divmod(i, 4)
        _kpi(s, x0 + c * (cw + gx), y0 + r * (ch + 0.2), cw, ch, value, label)
    # детализация — текстовый блок
    det = []
    op = int(tep.parking_open_places.value or 0)
    ml = int(tep.parking_multilevel_places.value or 0)
    ug = int(tep.parking_underground_places.value or 0)
    det.append(f"Парковки: всего {int(tep.parking_required_places.value or 0)} м/м  "
               f"(откр. {op} · многоур. {ml} · подз. {ug})")
    det.append(f"ДОО: {_objs_text(tep.kindergarten_places_accepted.formula, int(tep.kindergarten_places_accepted.value or 0))}  ·  "
               f"СОШ: {_objs_text(tep.school_places_accepted.formula, int(tep.school_places_accepted.value or 0))}")
    det.append(f"ЗНОП: {_fmt(tep.znop_area.value)} м² ({_fmt(tep.znop_per_person.value,1)} м²/чел)  ·  "
               f"Спорт: {_fmt(tep.sport_facilities_area.value)} м²")
    if tep.floor_clusters_detail:
        zones = " · ".join(
            f"{d['label']} {d['floors']}эт (КИТ {d['kit']:.2f})"
            for d in tep.floor_clusters_detail
        )
        det.append(f"Зоны этажности: {zones}")
    _text(s, 0.5, 4.55, 12.3, 2.5, "\n".join("•  " + d for d in det), size=14)

    # ─── Слайд 3: Баланс территории ───
    s = _slide(); _band(s, "Баланс территории", idx=2); _footer(s, name)
    chart = _chart_balance(tep)
    if chart is not None:
        # ширина ~8.0", чтобы не наезжать на правый текст (правее 8.6")
        s.shapes.add_picture(chart, int(0.4 * EMU), int(1.2 * EMU),
                             width=int(8.0 * EMU))
    b = tep.balance
    _text(s, 8.8, 1.4, 4.2, 0.5, "Ключевое", size=16, bold=True, color=_BLUE)
    _text(s, 8.8, 2.0, 4.2, 4.5, "\n".join("•  " + x for x in [
        f"Резерв (озеленение/двор): {_fmt(b.surplus)} м²",
        f"Занято всего: {_fmt(b.required_total)} м²",
        f"ЗУ жилой застройки: {_fmt(tep.housing_lot_area.value)} м²",
        f"Население: {_fmt(tep.population.value)} чел.",
        f"Плотность: {_fmt(tep.density_chel_per_ga.value,0)} чел/га",
    ]), size=13)

    # ─── Слайд 4: Жильё и соц. инфраструктура (таблица) ───
    s = _slide(); _band(s, "Жильё и социальная инфраструктура", idx=3); _footer(s, name)
    _table(s, 0.5, 1.3, 9.5,
           ["Показатель", "Значение", "Ед."],
           [
               ["Площадь квартир", _fmt(tep.apartments_area.value), "м²"],
               ["Общая площадь зданий (GFA)", _fmt(tep.gfa.value), "м²"],
               ["ВПП (встроенно-пристроенные)", _fmt(tep.built_in_area.value), "м²"],
               ["Население", _fmt(tep.population.value), "чел."],
               ["ДОО — требуется / принято",
                f"{_fmt(tep.kindergarten_places_required.value)} / "
                f"{_fmt(tep.kindergarten_places_accepted.value)}", "мест"],
               ["СОШ — требуется / принято",
                f"{_fmt(tep.school_places_required.value)} / "
                f"{_fmt(tep.school_places_accepted.value)}", "мест"],
               ["Озеленение жилого ЗУ", _fmt(tep.greening_housing_area.value), "м²"],
               ["ЗНОП", _fmt(tep.znop_area.value), "м²"],
           ],
           col_ratios=[3, 1.5, 0.8], fsize=12)

    # ─── Слайд 5: Парковки, проезды (+ кластеры) ───
    s = _slide(); _band(s, "Парковки и территории", idx=4); _footer(s, name)
    _table(s, 0.5, 1.25, 8.0,
           ["Парковки", "м/м", "Площадь, м²"],
           [
               ["Всего требуется", _fmt(tep.parking_required_places.value), "—"],
               ["Открытые наземные", _fmt(tep.parking_open_places.value),
                _fmt(tep.parking_open_area.value)],
               ["Многоуровневые", _fmt(tep.parking_multilevel_places.value),
                _fmt(tep.parking_multilevel_area.value)],
               ["Подземные", _fmt(tep.parking_underground_places.value), "—"],
               ["Соцобъектов (ДОО/СОШ)", _fmt(tep.social_parking_total.value),
                _fmt(tep.social_parking_area.value)],
           ],
           col_ratios=[2.4, 1, 1.4], fsize=12)
    if tep.floor_clusters_detail:
        _table(s, 0.5, 4.4, 9.5,
               ["Зона", "Площ., м²", "Этажей", "КИТ зоны", "Кварт., м²"],
               [[d["label"], _fmt(d["area_m2"]), str(d["floors"]),
                 _fmt(d["kit"], 3), _fmt(d["apartments_area"])]
                for d in tep.floor_clusters_detail],
               col_ratios=[1.3, 1.2, 0.9, 1, 1.2], fsize=11)

    # ─── Слайд 6 (последний): Экономика ───
    if tep.economy is not None:
        s = _slide(); _band(s, "Экономическая оценка", idx=5); _footer(s, name)
        ce = _chart_economy(tep)
        if ce is not None:
            s.shapes.add_picture(ce, int(0.4 * EMU), int(1.2 * EMU),
                                 width=int(7.2 * EMU))
        e = tep.economy
        for i, (v, lbl) in enumerate([
            (_fmt(e.profit), "Выгодность, баллы"),
            (_fmt(e.profit_before_social), "Без соц. нагрузки"),
            (f"{e.margin*100:.1f}%", "Маржа"),
            (f"{e.sellable_ratio*100:.0f}%", "Выход жилья"),
        ]):
            _kpi(s, 8.2, 1.4 + i * 1.32, 4.6, 1.1, v, lbl)
        _text(s, 0.4, 6.9, 12.5, 0.5,
              "«Баллы выгодности» — безразмерный индикатор сравнения вариантов "
              "(1.0 ≈ м² жилья 9-эт. монолита), не денежная оценка.",
              size=9, color=_GREY)

    prs.save(path)
    return path
